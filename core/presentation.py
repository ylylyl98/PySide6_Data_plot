"""PowerPoint assembly helpers for processed plot images.

The functions in this module are deliberately independent from Qt so slide
planning, safe output handling, and duplicate prevention can be tested without
starting the desktop interface.
"""

from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
import sys
import tempfile
from dataclasses import asdict, dataclass
from decimal import Decimal, InvalidOperation
from pathlib import Path
from typing import Iterable, Sequence


MANIFEST_SUFFIX = ".dptk.json"


@dataclass(frozen=True)
class PlotImage:
    path: Path
    relative_path: str
    workflow: str
    modified_time: float
    plot_kind: str = "other"


@dataclass(frozen=True)
class PresentationImage:
    path: Path
    caption: str = ""
    panel_label: str = ""
    logical_id: str = ""


@dataclass(frozen=True)
class GridLayout:
    rows: int
    columns: int


@dataclass(frozen=True)
class BuildResult:
    output_path: Path
    slides_added: int
    images_added: int
    images_skipped: int
    total_slides: int
    manifest_path: Path
    live_edit: bool = False
    saved: bool = True
    backup_path: Path | None = None
    images_updated: int = 0


@dataclass(frozen=True)
class PlannedSlide:
    images: tuple[PresentationImage, ...]
    group_key: str = ""
    group_label: str = ""
    part_index: int = 1
    part_count: int = 1


def plot_image_kind(path: str | Path, workflow: str = "") -> str:
    """Classify MCD exports for the Slides plot picker."""
    name = Path(path).name.casefold()
    is_mcd = str(workflow).casefold() == "mcd" or "_mcd_" in name
    if not is_mcd:
        return "other"
    if "_mcd_vs_b_" in name:
        return "mcd_b"
    if "_mcd_combo_" in name or "_mcd_combo." in name:
        return "mcd_combo"
    return "mcd_other"


def discover_plot_images(folder: str | Path) -> list[PlotImage]:
    """Return PNG plots below *folder*, newest first, with readable paths."""

    root = Path(folder).expanduser()
    if not root.is_dir():
        return []
    records: list[PlotImage] = []
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() != ".png":
            continue
        relative = path.relative_to(root)
        workflow = relative.parts[0] if len(relative.parts) > 1 else "Other"
        try:
            modified = path.stat().st_mtime
        except OSError:
            modified = 0.0
        records.append(
            PlotImage(
                path=path.resolve(),
                relative_path=str(relative),
                workflow=workflow,
                modified_time=modified,
                plot_kind=plot_image_kind(path, workflow),
            )
        )
    return sorted(records, key=lambda item: (-item.modified_time, item.relative_path.lower()))


def grid_for_count(count: int) -> GridLayout:
    """Choose a compact plot grid for one through twelve images."""

    if not 1 <= int(count) <= 12:
        raise ValueError("A slide supports between 1 and 12 images.")
    if count == 1:
        return GridLayout(1, 1)
    if count == 2:
        return GridLayout(1, 2)
    if count == 3:
        return GridLayout(1, 3)
    if count == 4:
        return GridLayout(2, 2)
    if count <= 6:
        return GridLayout(2, 3)
    if count <= 8:
        return GridLayout(2, 4)
    if count == 9:
        return GridLayout(3, 3)
    return GridLayout(3, 4)


def chunk_images(images: Sequence[PresentationImage], images_per_slide: int) -> list[list[PresentationImage]]:
    if not 1 <= int(images_per_slide) <= 12:
        raise ValueError("Images per slide must be between 1 and 12.")
    return [list(images[index:index + images_per_slide]) for index in range(0, len(images), images_per_slide)]


_NUMBER_TOKEN = r"[+-]?\d+(?:[p.]\d+)?"


def _number_from_token(token: str) -> Decimal | None:
    try:
        return Decimal(str(token).replace("p", "."))
    except InvalidOperation:
        return None


def _number_text(value: Decimal | None) -> str:
    if value is None:
        return "unknown"
    normalized = value.normalize()
    text = format(normalized, "f")
    return "0" if text in {"-0", "+0"} else text


def plot_metadata(path: str | Path) -> dict[str, object]:
    """Parse presentation-relevant metadata encoded in processed plot paths."""

    target = Path(path)
    text = "_".join((target.parent.name, target.stem))

    def number(pattern: str) -> Decimal | None:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        return _number_from_token(match.group(1)) if match else None

    doping = number(rf"(?:^|_)D({_NUMBER_TOKEN})(?:_|$)")
    vtg = number(rf"(?:^|_)Vtg({_NUMBER_TOKEN})(?:_|$)")
    vbg = number(rf"(?:^|_)Vbg({_NUMBER_TOKEN})(?:_|$)")
    efield = number(rf"(?:^|_)F({_NUMBER_TOKEN})(?:_|$)")
    energy = number(rf"(?:^|_)E({_NUMBER_TOKEN})eV(?:_|$)")
    window = number(rf"(?:^|_)W({_NUMBER_TOKEN})meV(?:_|$)")
    b_match = re.search(
        rf"(?:^|_)B({_NUMBER_TOKEN})to({_NUMBER_TOKEN})T(?:_|$)",
        text,
        flags=re.IGNORECASE,
    )
    b_range = (
        (_number_from_token(b_match.group(1)), _number_from_token(b_match.group(2)))
        if b_match else None
    )
    return {
        "doping": doping,
        "vtg": vtg,
        "vbg": vbg,
        "efield": efield,
        "energy": energy,
        "window": window,
        "b_range": b_range,
        "folder": target.parent.name,
        "plot_kind": plot_image_kind(target, "MCD" if "mcd" in text.casefold() else ""),
    }


def logical_result_id(path: str | Path) -> str:
    """Return a stable identity for one plotted result, independent of PNG bytes."""

    target = Path(path)
    metadata = plot_metadata(target)
    kind = str(metadata["plot_kind"])
    if kind in {"mcd_combo", "mcd_b"}:
        identity = {
            "dataset": target.parent.name.casefold(),
            "kind": kind,
            "energy_ev": _number_text(metadata["energy"]),
            "window_mev": _number_text(metadata["window"]),
        }
    else:
        identity = {
            "dataset": target.parent.name.casefold(),
            "kind": kind,
            "result": target.stem.casefold(),
        }
    payload = json.dumps(identity, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _group_identity(path: Path, group_by: str) -> tuple[str, str]:
    metadata = plot_metadata(path)
    mode = str(group_by or "queue").casefold()
    if mode == "doping":
        value = metadata["doping"]
        normalized = _number_text(value)
        return f"doping:{normalized}", f"Doping = {normalized} V"
    if mode == "folder":
        value = str(path.parent.resolve()).casefold()
        return f"folder:{value}", path.parent.name
    if mode == "gate":
        vtg = _number_text(metadata["vtg"])
        vbg = _number_text(metadata["vbg"])
        return f"gate:{vtg}:{vbg}", f"Vtg = {vtg} V, Vbg = {vbg} V"
    if mode == "efield":
        value = _number_text(metadata["efield"])
        return f"efield:{value}", f"E-field = {value} V"
    if mode == "b_range":
        values = metadata["b_range"]
        if values:
            start, end = (_number_text(item) for item in values)
            return f"b:{start}:{end}", f"B = {start}→{end} T"
        return "b:unknown", "B range unknown"
    return "queue", ""


def _mcd_order_key(image: PresentationImage, folder_order: dict[str, int], queue_index: int):
    path = Path(image.path)
    metadata = plot_metadata(path)
    folder = str(path.parent.resolve()).casefold()
    kind_order = {"mcd_combo": 0, "mcd_b": 1}.get(str(metadata["plot_kind"]), 2)
    energy = metadata["energy"]
    energy_key = float(energy) if isinstance(energy, Decimal) else float("inf")
    return folder_order[folder], kind_order, energy_key, queue_index


def plan_presentation_slides(
    images: Sequence[PresentationImage],
    images_per_slide: int,
    group_by: str = "queue",
) -> list[PlannedSlide]:
    """Group and order images consistently for preview, saved, and live output."""

    if not 1 <= int(images_per_slide) <= 12:
        raise ValueError("Images per slide must be between 1 and 12.")
    if not images:
        return []
    grouped: dict[str, tuple[str, list[tuple[int, PresentationImage]]]] = {}
    for index, image in enumerate(images):
        key, label = _group_identity(Path(image.path), group_by)
        grouped.setdefault(key, (label, []))[1].append((index, image))

    slides: list[PlannedSlide] = []
    for key, (label, indexed_images) in grouped.items():
        folder_order: dict[str, int] = {}
        for _index, image in indexed_images:
            folder = str(Path(image.path).parent.resolve()).casefold()
            folder_order.setdefault(folder, len(folder_order))
        ordered = [
            image
            for queue_index, image in sorted(
                indexed_images,
                key=lambda pair: _mcd_order_key(pair[1], folder_order, pair[0]),
            )
        ]
        chunks = chunk_images(ordered, images_per_slide)
        for part_index, chunk in enumerate(chunks, start=1):
            slides.append(
                PlannedSlide(
                    images=tuple(chunk),
                    group_key=key,
                    group_label=label,
                    part_index=part_index,
                    part_count=len(chunks),
                )
            )
    return slides


def planned_slide_title(plan: PlannedSlide, title_prefix: str = "") -> str:
    parts = [part for part in (title_prefix.strip(), plan.group_label) if part]
    title = " · ".join(parts)
    if plan.part_count > 1:
        suffix = f"{plan.part_index}/{plan.part_count}"
        return f"{title} ({suffix})" if title else ""
    return title


def compact_caption(path: str | Path, maximum: int = 64) -> str:
    """Create a short editable caption without changing the source image."""

    text = Path(path).stem.replace("_", " ")
    text = re.sub(r"\s+", " ", text).strip()
    text = re.sub(r"\b(?:YLin|YLog|linear|log)\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text).strip(" -_")
    if len(text) <= maximum:
        return text
    keep = max(10, (maximum - 3) // 2)
    return f"{text[:keep].rstrip()}...{text[-keep:].lstrip()}"


def panel_label(index: int) -> str:
    """Return spreadsheet-style panel labels: A..Z, AA..AZ, and so on."""

    if index < 0:
        raise ValueError("Panel index cannot be negative.")
    label = ""
    value = index
    while True:
        value, remainder = divmod(value, 26)
        label = chr(ord("A") + remainder) + label
        if value == 0:
            return label
        value -= 1


def file_sha256(path: str | Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def manifest_path_for(output_path: str | Path) -> Path:
    output = Path(output_path)
    return output.with_suffix(output.suffix + MANIFEST_SUFFIX)


def default_output_path(source_path: str | Path | None, image_root: str | Path | None = None) -> Path:
    if source_path:
        source = Path(source_path)
        return source.with_name(f"{source.stem}_with_plots.pptx")
    root = Path(image_root) if image_root else Path.cwd()
    return root / "DPTK_Presentation.pptx"


def unused_output_path(path: str | Path) -> Path:
    """Return *path* unless occupied by a non-DPTK file, then add a number."""

    candidate = Path(path)
    if not candidate.exists() or manifest_path_for(candidate).is_file():
        return candidate
    for number in range(2, 10000):
        numbered = candidate.with_name(f"{candidate.stem}_{number}{candidate.suffix}")
        if not numbered.exists():
            return numbered
    raise FileExistsError(f"Could not find an unused output name near {candidate}.")


def _load_manifest(path: Path) -> dict:
    if not path.is_file():
        return {"format": 1, "images": [], "builds": []}
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {"format": 1, "images": [], "builds": []}
    if not isinstance(payload, dict):
        return {"format": 1, "images": [], "builds": []}
    payload.setdefault("format", 1)
    payload.setdefault("images", [])
    payload.setdefault("builds", [])
    return payload


def _fit_size(image_path: Path, box_width: int, box_height: int) -> tuple[int, int]:
    from PIL import Image

    with Image.open(image_path) as image:
        width, height = image.size
    if width <= 0 or height <= 0:
        raise ValueError(f"Invalid image dimensions: {image_path}")
    scale = min(box_width / width, box_height / height)
    return max(1, int(width * scale)), max(1, int(height * scale))


def _add_text_box(slide, text: str, left: int, top: int, width: int, height: int, *, points: float, bold: bool = False, centered: bool = False) -> None:
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
    paragraph.font.size = Pt(points)
    paragraph.font.bold = bold


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _pptx_existing_logical_ids(prs) -> set[str]:
    prefix = "DPTK_RESULT_"
    result: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            name = str(getattr(shape, "name", "") or "")
            if name.startswith(prefix):
                result_id = name[len(prefix):].strip()
                if result_id:
                    result.add(result_id)
    return result


def _add_plot_slide(
    prs,
    images: Sequence[PresentationImage],
    title: str,
    *,
    group_key: str = "",
    show_captions: bool,
    show_panel_labels: bool,
) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(_blank_layout(prs))
    if group_key:
        slide._element.cSld.set("name", f"DPTK_GROUP_{group_key}")
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    outer = int(Inches(0.28))
    gap = int(Inches(0.12))
    title_height = int(Inches(0.5)) if title else 0
    title_gap = int(Inches(0.08)) if title else 0
    if title:
        _add_text_box(
            slide,
            title,
            outer,
            int(Inches(0.08)),
            slide_width - 2 * outer,
            title_height,
            points=24,
            bold=True,
        )

    layout = grid_for_count(len(images))
    body_top = outer + title_height + title_gap
    body_height = slide_height - body_top - outer
    cell_width = (slide_width - 2 * outer - (layout.columns - 1) * gap) // layout.columns
    cell_height = (body_height - (layout.rows - 1) * gap) // layout.rows
    caption_height = int(Inches(0.28)) if show_captions else 0

    for index, image in enumerate(images):
        row, column = divmod(index, layout.columns)
        cell_left = outer + column * (cell_width + gap)
        cell_top = body_top + row * (cell_height + gap)
        picture_height = max(1, cell_height - caption_height)
        fitted_width, fitted_height = _fit_size(image.path, cell_width, picture_height)
        picture_left = cell_left + (cell_width - fitted_width) // 2
        picture_top = cell_top + (picture_height - fitted_height) // 2
        picture = slide.shapes.add_picture(
            str(image.path), picture_left, picture_top, fitted_width, fitted_height
        )
        result_id = image.logical_id or logical_result_id(image.path)
        picture.name = f"DPTK_RESULT_{result_id}"
        if show_panel_labels and image.panel_label:
            _add_text_box(
                slide,
                image.panel_label,
                cell_left + int(Inches(0.04)),
                cell_top + int(Inches(0.02)),
                int(Inches(0.3)),
                int(Inches(0.22)),
                points=11,
                bold=True,
                centered=True,
            )
        if show_captions and image.caption:
            _add_text_box(
                slide,
                image.caption,
                cell_left,
                cell_top + cell_height - caption_height,
                cell_width,
                caption_height,
                points=9 if len(images) <= 6 else 8,
                centered=True,
            )


def build_presentation(
    images: Sequence[PresentationImage],
    output_path: str | Path,
    *,
    source_path: str | Path | None = None,
    images_per_slide: int = 6,
    title_prefix: str = "",
    show_captions: bool = True,
    show_panel_labels: bool = True,
    group_by: str = "queue",
    in_place: bool = False,
    create_backup: bool = True,
) -> BuildResult:
    """Append plot slides to a deck copy or atomically update the source deck."""

    from pptx import Presentation

    if not images:
        raise ValueError("Add at least one PNG to the presentation queue.")
    source = Path(source_path).expanduser().resolve() if source_path else None
    if source and (not source.is_file() or source.suffix.lower() != ".pptx"):
        raise FileNotFoundError(f"PowerPoint source not found: {source}")
    if in_place and source is None:
        raise ValueError("Choose an existing PowerPoint before inserting slides in place.")
    output = source if in_place else Path(output_path).expanduser().resolve()
    assert output is not None
    if output.suffix.lower() != ".pptx":
        output = output.with_suffix(".pptx")
    if source and output == source and not in_place:
        output = default_output_path(source)
    if not in_place:
        output = unused_output_path(output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    manifest_path = manifest_path_for(output)
    manifest = _load_manifest(manifest_path)

    continuing = output.is_file() and manifest_path.is_file()
    if not continuing:
        # A stale sidecar must never suppress plots when its presentation was
        # deleted or when we are starting from a different untouched deck.
        manifest = {"format": 1, "images": [], "builds": []}
    elif source:
        recorded_source = manifest.get("source_presentation")
        if not recorded_source or Path(str(recorded_source)).expanduser().resolve() != source:
            raise ValueError(
                "This output presentation belongs to a different source deck. "
                "Choose a new output filename before building."
            )
    deck_input = output if continuing or in_place else source
    prs = Presentation(str(deck_input)) if deck_input else Presentation()
    starting_slide_count = len(prs.slides)
    existing_hashes = {
        str(item.get("sha256"))
        for item in manifest.get("images", [])
        if isinstance(item, dict) and item.get("sha256")
    }
    existing_logical_ids = {
        str(item.get("logical_id") or logical_result_id(item.get("path")))
        for item in manifest.get("images", [])
        if isinstance(item, dict) and item.get("path")
    }
    existing_logical_ids.update(_pptx_existing_logical_ids(prs))
    accepted: list[tuple[PresentationImage, str, str]] = []
    skipped = 0
    for image in images:
        path = Path(image.path).expanduser().resolve()
        if not path.is_file() or path.suffix.lower() != ".png":
            raise FileNotFoundError(f"PNG not found: {path}")
        digest = file_sha256(path)
        result_id = image.logical_id or logical_result_id(path)
        if digest in existing_hashes or result_id in existing_logical_ids:
            skipped += 1
            continue
        accepted.append(
            (PresentationImage(path, image.caption, image.panel_label, result_id), digest, result_id)
        )
        existing_hashes.add(digest)
        existing_logical_ids.add(result_id)

    plans = plan_presentation_slides(
        [item for item, _digest, _result_id in accepted], images_per_slide, group_by
    )
    digest_by_id = {result_id: digest for _item, digest, result_id in accepted}
    added_records: list[dict] = []
    for plan in plans:
        group = list(plan.images)
        title = planned_slide_title(plan, title_prefix)
        _add_plot_slide(
            prs,
            group,
            title,
            group_key=plan.group_key,
            show_captions=show_captions,
            show_panel_labels=show_panel_labels,
        )
        slide_number = len(prs.slides)
        layout = grid_for_count(len(group))
        for item in group:
            added_records.append(
                {
                    "path": str(item.path),
                    "sha256": digest_by_id[item.logical_id],
                    "logical_id": item.logical_id,
                    "caption": item.caption,
                    "panel_label": item.panel_label,
                    "slide_number": slide_number,
                    "grid": asdict(layout),
                    "group_by": group_by,
                    "group_key": plan.group_key,
                }
            )

    backup_path: Path | None = None
    if in_place and create_backup and accepted:
        backup_path = output.with_name(f"{output.stem}.dptk-backup{output.suffix}")
        if not backup_path.exists():
            shutil.copy2(output, backup_path)

    if accepted or not output.exists():
        file_descriptor, temporary_name = tempfile.mkstemp(
            prefix=f".{output.stem}_", suffix=".pptx", dir=str(output.parent)
        )
        os.close(file_descriptor)
        temporary = Path(temporary_name)
        try:
            prs.save(str(temporary))
            os.replace(temporary, output)
        finally:
            if temporary.exists():
                temporary.unlink()

    manifest["format"] = 2
    manifest["source_presentation"] = str(source) if source else None
    manifest["output_presentation"] = str(output)
    manifest.setdefault("images", []).extend(added_records)
    manifest.setdefault("builds", []).append(
        {
            "slides_added": len(plans),
            "images_added": len(accepted),
            "images_skipped": skipped,
            "images_per_slide": images_per_slide,
            "group_by": group_by,
        }
    )
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return BuildResult(
        output_path=output,
        slides_added=len(plans),
        images_added=len(accepted),
        images_skipped=skipped,
        total_slides=starting_slide_count + len(plans),
        manifest_path=manifest_path,
        backup_path=backup_path,
    )


def _powerpoint_modules():
    if sys.platform != "win32":
        raise RuntimeError("Live PowerPoint editing is available in the Windows DPTK application.")
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise RuntimeError(
            "This DPTK build does not include live PowerPoint integration. "
            "Install the latest DPTK Windows build; no separate Python package is required."
        ) from exc
    return pythoncom, win32com.client


def powerpoint_integration_available() -> tuple[bool, str]:
    """Report whether this build contains the Windows PowerPoint bridge."""
    try:
        _powerpoint_modules()
    except RuntimeError as exc:
        return False, str(exc)
    return True, "Live PowerPoint integration is available."


def _open_powerpoint_presentation(application, presentation_path: Path):
    target = os.path.normcase(str(presentation_path.resolve()))
    for index in range(1, int(application.Presentations.Count) + 1):
        presentation = application.Presentations.Item(index)
        try:
            candidate = os.path.normcase(str(Path(str(presentation.FullName)).resolve()))
        except (OSError, ValueError):
            continue
        if candidate == target:
            return presentation
    return None


def powerpoint_presentation_is_open(presentation_path: str | Path) -> bool:
    """Return whether the exact saved deck is open in desktop PowerPoint."""
    path = Path(presentation_path).expanduser().resolve()
    try:
        pythoncom, win32_client = _powerpoint_modules()
    except RuntimeError:
        return False
    pythoncom.CoInitialize()
    try:
        try:
            application = win32_client.GetActiveObject("PowerPoint.Application")
        except Exception:
            return False
        return _open_powerpoint_presentation(application, path) is not None
    finally:
        pythoncom.CoUninitialize()


def _live_existing_plot_ids(presentation) -> tuple[set[str], set[str]]:
    hashes: set[str] = set()
    logical_ids: set[str] = set()
    for slide_index in range(1, int(presentation.Slides.Count) + 1):
        slide = presentation.Slides.Item(slide_index)
        for shape_index in range(1, int(slide.Shapes.Count) + 1):
            shape = slide.Shapes.Item(shape_index)
            try:
                value = str(shape.Tags.Item("DPTK_SHA256") or "").strip()
            except Exception:
                continue
            if value:
                hashes.add(value)
            try:
                result_id = str(shape.Tags.Item("DPTK_RESULT_ID") or "").strip()
            except Exception:
                result_id = ""
            if result_id:
                logical_ids.add(result_id)
    return hashes, logical_ids


def _add_live_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    points: float,
    bold: bool = False,
    centered: bool = False,
):
    shape = slide.Shapes.AddTextbox(1, left, top, width, height)
    frame = shape.TextFrame
    frame.MarginLeft = frame.MarginRight = 0
    frame.MarginTop = frame.MarginBottom = 0
    frame.VerticalAnchor = 3
    text_range = frame.TextRange
    text_range.Text = text
    text_range.Font.Size = float(points)
    text_range.Font.Bold = -1 if bold else 0
    text_range.ParagraphFormat.Alignment = 2 if centered else 1
    return shape


def _add_live_plot_slide(
    presentation,
    images: Sequence[tuple[PresentationImage, str, str]],
    title: str,
    *,
    group_key: str = "",
    show_captions: bool,
    show_panel_labels: bool,
) -> tuple[object, GridLayout]:
    slide = presentation.Slides.Add(int(presentation.Slides.Count) + 1, 12)
    if group_key:
        try:
            slide.Tags.Add("DPTK_GROUP_KEY", group_key)
        except Exception:
            pass
    slide_width = float(presentation.PageSetup.SlideWidth)
    slide_height = float(presentation.PageSetup.SlideHeight)
    outer = 0.28 * 72.0
    gap = 0.12 * 72.0
    title_height = 0.5 * 72.0 if title else 0.0
    title_gap = 0.08 * 72.0 if title else 0.0
    if title:
        _add_live_text_box(
            slide, title, outer, 0.08 * 72.0, slide_width - 2 * outer,
            title_height, points=24, bold=True,
        )
    layout = grid_for_count(len(images))
    body_top = outer + title_height + title_gap
    body_height = slide_height - body_top - outer
    cell_width = (slide_width - 2 * outer - (layout.columns - 1) * gap) / layout.columns
    cell_height = (body_height - (layout.rows - 1) * gap) / layout.rows
    caption_height = 0.28 * 72.0 if show_captions else 0.0
    for index, (image, digest, result_id) in enumerate(images):
        row, column = divmod(index, layout.columns)
        cell_left = outer + column * (cell_width + gap)
        cell_top = body_top + row * (cell_height + gap)
        picture_height = max(1.0, cell_height - caption_height)
        fitted_width, fitted_height = _fit_size(
            image.path, int(cell_width), int(picture_height)
        )
        picture_left = cell_left + (cell_width - fitted_width) / 2.0
        picture_top = cell_top + (picture_height - fitted_height) / 2.0
        picture = slide.Shapes.AddPicture(
            str(image.path), 0, -1, picture_left, picture_top,
            float(fitted_width), float(fitted_height),
        )
        picture.Tags.Add("DPTK_SHA256", digest)
        picture.Tags.Add("DPTK_RESULT_ID", result_id)
        if show_panel_labels and image.panel_label:
            _add_live_text_box(
                slide, image.panel_label, cell_left + 0.04 * 72.0,
                cell_top + 0.02 * 72.0, 0.3 * 72.0, 0.22 * 72.0,
                points=11, bold=True, centered=True,
            )
        if show_captions and image.caption:
            _add_live_text_box(
                slide, image.caption, cell_left,
                cell_top + cell_height - caption_height,
                cell_width, caption_height,
                points=9 if len(images) <= 6 else 8, centered=True,
            )
    return slide, layout


def insert_plots_into_open_powerpoint(
    images: Sequence[PresentationImage],
    presentation_path: str | Path,
    *,
    images_per_slide: int = 6,
    title_prefix: str = "",
    show_captions: bool = True,
    show_panel_labels: bool = True,
    group_by: str = "queue",
    save: bool = False,
    create_backup: bool = True,
) -> BuildResult:
    """Insert plots into the exact open PowerPoint deck through COM."""
    if not images:
        raise ValueError("Add at least one PNG to the presentation queue.")
    target = Path(presentation_path).expanduser().resolve()
    if not target.is_file() or target.suffix.lower() != ".pptx":
        raise FileNotFoundError(f"PowerPoint presentation not found: {target}")
    pythoncom, win32_client = _powerpoint_modules()
    pythoncom.CoInitialize()
    try:
        try:
            application = win32_client.GetActiveObject("PowerPoint.Application")
        except Exception as exc:
            raise RuntimeError("PowerPoint is not currently open.") from exc
        presentation = _open_powerpoint_presentation(application, target)
        if presentation is None:
            raise RuntimeError(
                f"The selected presentation is not open in PowerPoint: {target.name}"
            )
        if bool(getattr(presentation, "ReadOnly", False)):
            raise PermissionError("The selected PowerPoint presentation is read-only.")

        manifest_path = manifest_path_for(target)
        manifest = _load_manifest(manifest_path)
        existing_hashes = {
            str(item.get("sha256"))
            for item in manifest.get("images", [])
            if isinstance(item, dict) and item.get("sha256")
        }
        manifest_logical_ids = {
            str(item.get("logical_id") or logical_result_id(item.get("path")))
            for item in manifest.get("images", [])
            if isinstance(item, dict) and item.get("path")
        }
        live_hashes, live_logical_ids = _live_existing_plot_ids(presentation)
        existing_hashes.update(live_hashes)
        existing_logical_ids = manifest_logical_ids | live_logical_ids
        accepted: list[tuple[PresentationImage, str, str]] = []
        skipped = 0
        for image in images:
            path = Path(image.path).expanduser().resolve()
            if not path.is_file() or path.suffix.lower() != ".png":
                raise FileNotFoundError(f"PNG not found: {path}")
            digest = file_sha256(path)
            result_id = image.logical_id or logical_result_id(path)
            if digest in existing_hashes or result_id in existing_logical_ids:
                skipped += 1
                continue
            accepted.append(
                (PresentationImage(path, image.caption, image.panel_label, result_id), digest, result_id)
            )
            existing_hashes.add(digest)
            existing_logical_ids.add(result_id)

        backup_path: Path | None = None
        if create_backup and accepted:
            backup_path = target.with_name(f"{target.stem}.dptk-backup{target.suffix}")
            if not backup_path.exists():
                presentation.SaveCopyAs(str(backup_path))

        plans = plan_presentation_slides(
            [item for item, _digest, _result_id in accepted], images_per_slide, group_by
        )
        accepted_by_id = {
            result_id: (item, digest, result_id)
            for item, digest, result_id in accepted
        }
        starting_slide_count = int(presentation.Slides.Count)
        added_records: list[dict] = []
        for plan in plans:
            group = list(plan.images)
            title = planned_slide_title(plan, title_prefix)
            tagged_group = [accepted_by_id[item.logical_id] for item in group]
            slide, layout = _add_live_plot_slide(
                presentation, tagged_group, title,
                group_key=plan.group_key,
                show_captions=show_captions,
                show_panel_labels=show_panel_labels,
            )
            slide_number = int(slide.SlideIndex)
            for item, digest, result_id in tagged_group:
                added_records.append({
                    "path": str(item.path), "sha256": digest,
                    "logical_id": result_id,
                    "caption": item.caption, "panel_label": item.panel_label,
                    "slide_number": slide_number, "grid": asdict(layout),
                    "group_by": group_by, "group_key": plan.group_key,
                })
        if plans:
            try:
                application.ActiveWindow.View.GotoSlide(int(presentation.Slides.Count))
            except Exception:
                # Insertion still succeeded when a different PowerPoint window
                # is active or the deck is not in Normal view.
                pass
        if save and accepted:
            presentation.Save()
            manifest["format"] = 2
            manifest["source_presentation"] = str(target)
            manifest["output_presentation"] = str(target)
            manifest.setdefault("images", []).extend(added_records)
            manifest.setdefault("builds", []).append({
                "slides_added": len(plans), "images_added": len(accepted),
                "images_skipped": skipped, "images_per_slide": images_per_slide,
                "group_by": group_by, "live_edit": True,
            })
            manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
        return BuildResult(
            output_path=target,
            slides_added=len(plans), images_added=len(accepted),
            images_skipped=skipped,
            total_slides=starting_slide_count + len(plans),
            manifest_path=manifest_path,
            live_edit=True, saved=bool(save), backup_path=backup_path,
        )
    finally:
        pythoncom.CoUninitialize()
