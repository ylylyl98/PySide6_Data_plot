from __future__ import annotations

import json
import os
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.presentation import (
    _add_live_plot_slide,
    PresentationImage,
    build_presentation,
    compact_caption,
    default_output_path,
    discover_plot_images,
    file_sha256,
    grid_for_count,
    logical_result_id,
    manifest_path_for,
    panel_label,
    plan_presentation_slides,
    planned_slide_title,
    plot_metadata,
    plot_image_kind,
)


class PresentationTests(unittest.TestCase):
    def _png(self, folder: Path, name: str, size=(640, 360), color=(30, 90, 180)) -> Path:
        path = folder / name
        path.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", size, color).save(path)
        return path

    def test_grid_layouts_cover_every_supported_count(self) -> None:
        expected = {
            1: (1, 1),
            2: (1, 2),
            3: (1, 3),
            4: (2, 2),
            5: (2, 3),
            6: (2, 3),
            7: (2, 4),
            8: (2, 4),
            9: (3, 3),
            10: (3, 4),
            11: (3, 4),
            12: (3, 4),
        }
        for count, dimensions in expected.items():
            layout = grid_for_count(count)
            self.assertEqual((layout.rows, layout.columns), dimensions)
        with self.assertRaises(ValueError):
            grid_for_count(0)
        with self.assertRaises(ValueError):
            grid_for_count(13)

    def test_discovery_uses_full_relative_paths_and_newest_first(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            first = self._png(root, "DRR/group_with_a_long_name/plot_one.png")
            second = self._png(root, "PL/another_group/plot_two.PNG")
            os.utime(second, (100, 100))
            os.utime(first, (200, 200))
            records = discover_plot_images(root)
            self.assertEqual({record.path for record in records}, {first.resolve(), second.resolve()})
            self.assertIn("group_with_a_long_name", records[0].relative_path)
            self.assertEqual({record.workflow for record in records}, {"DRR", "PL"})

    def test_discovery_classifies_mcd_combo_and_field_trace(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            combo = self._png(root, "MCD/run/sample_MCD_Combo_map.png")
            trace = self._png(root, "MCD/run/sample_MCD_vs_B_E1p64.png")
            other = self._png(root, "DRR/run/sample_DRR.png")
            kinds = {record.path: record.plot_kind for record in discover_plot_images(root)}
            self.assertEqual(kinds[combo.resolve()], "mcd_combo")
            self.assertEqual(kinds[trace.resolve()], "mcd_b")
            self.assertEqual(kinds[other.resolve()], "other")
            self.assertEqual(plot_image_kind("unknown_MCD_plot.png", "MCD"), "mcd_other")

    def test_captions_and_panel_labels_are_readable(self) -> None:
        caption = compact_caption("sample_DRR_YLin_a_very_long_measurement_filename_that_needs_shortening.png", 36)
        self.assertNotIn("YLin", caption)
        self.assertLessEqual(len(caption), 37)
        self.assertEqual([panel_label(index) for index in (0, 1, 25, 26, 27)], ["A", "B", "Z", "AA", "AB"])

    def test_metadata_normalizes_folder_values_and_plans_repeated_doping_together(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            paths = [
                self._png(root, "MCD/run_D6p30_F0_B-2to+2T/run_MCD_vs_B_E1.640000eV_W5meV.png"),
                self._png(root, "MCD/run_D2_F0_B-2to+2T/run_MCD_vs_B_E1.60eV_W5meV.png"),
                self._png(root, "MCD/repeat_D6p3_F20_B-9to+9T/run_MCD_Combo.png"),
                self._png(root, "MCD/repeat_D6p3_F20_B-9to+9T/run_MCD_vs_B_E1.58eV_W5meV.png"),
            ]
            metadata = plot_metadata(paths[0])
            self.assertEqual(str(metadata["doping"]), "6.30")
            self.assertEqual(tuple(str(value) for value in metadata["b_range"]), ("-2", "2"))

            plans = plan_presentation_slides(
                [PresentationImage(path) for path in paths], 12, "doping"
            )
            self.assertEqual([len(plan.images) for plan in plans], [3, 1])
            self.assertEqual(plans[0].group_key, "doping:6.3")
            self.assertEqual(Path(plans[0].images[1].path), paths[2])
            self.assertEqual(Path(plans[0].images[2].path), paths[3])
            self.assertEqual(planned_slide_title(plans[0]), "Doping = 6.3 V")

    def test_exported_panel_labels_restart_after_automatic_grouping(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            paths = [
                self._png(root, "MCD/run_D1_F0/run_D1_F0_MCD_vs_B_E1p57eV_W5meV.png", color=(10, 20, 30)),
                self._png(root, "MCD/run_D2_F0/run_D2_F0_MCD_vs_B_E1p60eV_W5meV.png", color=(40, 50, 60)),
                self._png(root, "MCD/run_D1_F20/run_D1_F20_MCD_vs_B_E1p64eV_W5meV.png", color=(70, 80, 90)),
            ]
            output = root / "grouped.pptx"
            build_presentation(
                [PresentationImage(path) for path in paths],
                output,
                images_per_slide=12,
                group_by="doping",
                show_panel_labels=True,
            )

            manifest = json.loads(manifest_path_for(output).read_text(encoding="utf-8"))
            self.assertEqual(
                [item["panel_label"] for item in manifest["images"]],
                ["A", "B", "A"],
            )

    def test_live_insert_names_picture_for_closed_deck_duplicate_detection(self) -> None:
        class Tags:
            def Add(self, _name, _value) -> None:
                pass

        class Picture:
            def __init__(self) -> None:
                self.Name = "Picture 1"
                self.Tags = Tags()

        class Shapes:
            def __init__(self) -> None:
                self.picture = Picture()

            def AddPicture(self, *_args):
                return self.picture

        class Slide:
            def __init__(self) -> None:
                self.Shapes = Shapes()
                self.Tags = Tags()

        class Slides:
            Count = 0

            def __init__(self) -> None:
                self.slide = Slide()

            def Add(self, *_args):
                return self.slide

        class PageSetup:
            SlideWidth = 960.0
            SlideHeight = 540.0

        class PresentationStub:
            def __init__(self) -> None:
                self.Slides = Slides()
                self.PageSetup = PageSetup()

        with tempfile.TemporaryDirectory() as tmp:
            image = self._png(Path(tmp), "plot.png")
            presentation = PresentationStub()
            _add_live_plot_slide(
                presentation,
                [(PresentationImage(image), "digest", "logical-id")],
                "",
                show_captions=False,
                show_panel_labels=False,
            )
            self.assertEqual(
                presentation.Slides.slide.Shapes.picture.Name,
                "DPTK_RESULT_logical-id",
            )

    def test_logical_identity_normalizes_mcd_energy_precision(self) -> None:
        first = Path("run_D1") / "sample_MCD_vs_B_E1.640000eV_W5meV.png"
        second = Path("run_D1") / "renamed_MCD_vs_B_E1p64eV_W5p0meV.png"
        self.assertEqual(logical_result_id(first), logical_result_id(second))

    def test_build_preserves_template_pngs_and_suppresses_repeat_duplicates(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "research_template.pptx"
            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[0])
            slide.shapes.title.text = "Existing research slide"
            template.save(source)

            paths = [
                self._png(root, f"Processed Data/DRR/plot_{index}.png", size=(900, 420 + index * 20), color=(20 + index, 80, 160))
                for index in range(7)
            ]
            hashes_before = {path: file_sha256(path) for path in paths}
            images = [
                PresentationImage(path, caption=f"Plot {index + 1}", panel_label=panel_label(index % 6))
                for index, path in enumerate(paths)
            ]
            output = default_output_path(source)
            result = build_presentation(
                images,
                output,
                source_path=source,
                images_per_slide=6,
                title_prefix="DRR comparison",
            )
            self.assertEqual(result.slides_added, 2)
            self.assertEqual(result.images_added, 7)
            self.assertEqual(result.total_slides, 3)
            self.assertTrue(output.is_file())
            self.assertEqual(len(Presentation(source).slides), 1)

            built = Presentation(output)
            self.assertEqual(len(built.slides), 3)
            self.assertEqual(built.slides[0].shapes.title.text, "Existing research slide")
            pictures = [
                shape
                for added_slide in list(built.slides)[1:]
                for shape in added_slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            self.assertEqual(len(pictures), 7)
            for picture in pictures:
                self.assertEqual((picture.crop_left, picture.crop_right, picture.crop_top, picture.crop_bottom), (0.0, 0.0, 0.0, 0.0))
            self.assertEqual(hashes_before, {path: file_sha256(path) for path in paths})

            repeated = build_presentation(
                images,
                output,
                source_path=source,
                images_per_slide=6,
                title_prefix="DRR comparison",
            )
            self.assertEqual(repeated.slides_added, 0)
            self.assertEqual(repeated.images_added, 0)
            self.assertEqual(repeated.images_skipped, 7)
            self.assertEqual(len(Presentation(output).slides), 3)
            manifest = json.loads(manifest_path_for(output).read_text(encoding="utf-8"))
            self.assertEqual(len(manifest["images"]), 7)
            self.assertEqual(len(manifest["builds"]), 2)

    def test_source_is_never_overwritten_when_used_as_output(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "deck.pptx"
            Presentation().save(source)
            image = self._png(root, "plot.png")
            source_hash = file_sha256(source)
            result = build_presentation(
                [PresentationImage(image, "Plot", "A")],
                source,
                source_path=source,
                images_per_slide=1,
            )
            self.assertNotEqual(result.output_path, source)
            self.assertEqual(file_sha256(source), source_hash)
            self.assertTrue(result.output_path.is_file())

    def test_regenerated_logical_result_is_not_inserted_twice(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image = self._png(
                root,
                "Processed Data/MCD/run_D1/run_D1_MCD_vs_B_E1.64eV_W5meV.png",
                color=(10, 20, 30),
            )
            output = root / "deck.pptx"
            first = build_presentation(
                [PresentationImage(image)], output, group_by="doping", images_per_slide=12
            )
            first_hash = file_sha256(image)
            Image.new("RGB", (640, 360), (200, 100, 50)).save(image)
            self.assertNotEqual(file_sha256(image), first_hash)
            repeated = build_presentation(
                [PresentationImage(image)], output, group_by="doping", images_per_slide=12
            )
            self.assertEqual(first.images_added, 1)
            self.assertEqual(repeated.images_added, 0)
            self.assertEqual(repeated.images_skipped, 1)
            self.assertEqual(len(Presentation(output).slides), 1)
            manifest = json.loads(manifest_path_for(output).read_text(encoding="utf-8"))
            self.assertTrue(manifest["images"][0]["logical_id"])
            self.assertEqual(manifest["images"][0]["group_key"], "doping:1")

    def test_in_place_deck_remembers_logical_results_without_sidecar(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "deck.pptx"
            Presentation().save(deck)
            image = self._png(
                root, "MCD/run_D1/run_D1_MCD_vs_B_E1.64eV_W5meV.png"
            )
            first = build_presentation(
                [PresentationImage(image)],
                deck,
                source_path=deck,
                in_place=True,
                group_by="doping",
            )
            manifest_path_for(deck).unlink()
            Image.new("RGB", (640, 360), (220, 40, 80)).save(image)
            repeated = build_presentation(
                [PresentationImage(image)],
                deck,
                source_path=deck,
                in_place=True,
                group_by="doping",
            )
            self.assertEqual(first.images_added, 1)
            self.assertEqual(repeated.images_added, 0)
            self.assertEqual(repeated.images_skipped, 1)
            self.assertEqual(len(Presentation(deck).slides), 1)

    def test_in_place_build_updates_selected_deck_and_creates_recovery_backup(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "deck_to_edit.pptx"
            deck = Presentation()
            deck.slides.add_slide(deck.slide_layouts[0])
            deck.save(source)
            image = self._png(root, "plot.png")
            result = build_presentation(
                [PresentationImage(image, "Plot", "A")],
                source,
                source_path=source,
                images_per_slide=1,
                in_place=True,
                create_backup=True,
            )
            self.assertEqual(result.output_path, source.resolve())
            self.assertEqual(len(Presentation(source).slides), 2)
            self.assertIsNotNone(result.backup_path)
            self.assertTrue(result.backup_path.is_file())
            self.assertEqual(len(Presentation(result.backup_path).slides), 1)

    def test_existing_output_cannot_be_reused_with_a_different_source_deck(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            first_source = root / "first.pptx"
            second_source = root / "second.pptx"
            Presentation().save(first_source)
            Presentation().save(second_source)
            image = self._png(root, "plot.png")
            output = root / "assembled.pptx"
            build_presentation(
                [PresentationImage(image, "Plot", "A")],
                output,
                source_path=first_source,
            )
            with self.assertRaisesRegex(ValueError, "different source deck"):
                build_presentation(
                    [PresentationImage(image, "Plot", "A")],
                    output,
                    source_path=second_source,
                )


if __name__ == "__main__":
    unittest.main()
